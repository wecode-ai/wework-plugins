#!/usr/bin/env python3
# SPDX-License-Identifier: Apache-2.0
"""Deterministic PPTX creation, inspection, editing, rendering, and validation."""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Inches, Pt

DEFAULT_THEME = {
    "background": "F8FAFC",
    "foreground": "0F172A",
    "muted": "475569",
    "accent": "2563EB",
}


def emit(payload: object, destination: str | None = None) -> None:
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    if destination:
        path = Path(destination)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(f"{text}\n", encoding="utf-8")
    else:
        print(text)


def read_spec(path: str) -> dict[str, object]:
    value = json.loads(Path(path).read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise TypeError("presentation specification must be a JSON object")
    return value


def rgb(value: object) -> RGBColor:
    normalized = str(value).strip().lstrip("#")
    if len(normalized) != 6:
        raise ValueError(f"invalid RGB color: {value}")
    return RGBColor.from_string(normalized.upper())


def theme_from(spec: dict[str, object]) -> dict[str, str]:
    theme = dict(DEFAULT_THEME)
    supplied = spec.get("theme", {})
    if isinstance(supplied, dict):
        for key in theme:
            if key in supplied:
                theme[key] = str(supplied[key]).lstrip("#")
    return theme


def add_textbox(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> object:
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def set_background(slide: Slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_footer(slide: Slide, index: int, theme: dict[str, str]) -> None:
    add_textbox(
        slide,
        str(index),
        12.45,
        7.02,
        0.4,
        0.24,
        size=8,
        color=theme["muted"],
        align=PP_ALIGN.RIGHT,
    )


def normalize_bullets(value: object) -> list[tuple[str, int]]:
    if not isinstance(value, list):
        raise TypeError("bullets must be an array")
    bullets: list[tuple[str, int]] = []
    for item in value:
        if isinstance(item, dict):
            bullets.append((str(item.get("text", "")), int(item.get("level", 0))))
        else:
            bullets.append((str(item), 0))
    return bullets


def add_bullets(
    slide: Slide,
    bullets: object,
    x: float,
    y: float,
    width: float,
    height: float,
    theme: dict[str, str],
) -> None:
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.04)
    for index, (text, level) in enumerate(normalize_bullets(bullets)):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.level = max(0, min(level, 4))
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(22 - min(level, 2) * 2)
        paragraph.font.color.rgb = rgb(theme["foreground"])
        paragraph.space_after = Pt(11)


def create_title_slide(
    presentation: PresentationType,
    value: dict[str, object],
    theme: dict[str, str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["background"])
    add_textbox(
        slide,
        str(value.get("title", "Untitled presentation")),
        0.9,
        1.55,
        11.5,
        1.6,
        size=36,
        color=theme["foreground"],
        bold=True,
    )
    subtitle = str(value.get("subtitle", "")).strip()
    if subtitle:
        add_textbox(
            slide,
            subtitle,
            0.95,
            3.25,
            10.8,
            0.8,
            size=19,
            color=theme["muted"],
        )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.95),
        Inches(4.45),
        Inches(1.3),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(theme["accent"])
    accent.line.fill.background()


def create_content_slide(
    presentation: PresentationType,
    value: dict[str, object],
    theme: dict[str, str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["background"])
    add_textbox(
        slide,
        str(value.get("title", "")),
        0.75,
        0.45,
        11.8,
        0.75,
        size=28,
        color=theme["foreground"],
        bold=True,
    )
    columns = value.get("columns")
    if isinstance(columns, list) and len(columns) == 2:
        for column_index, column in enumerate(columns):
            if not isinstance(column, dict):
                raise TypeError("column entries must be objects")
            x = 0.8 + column_index * 6.15
            heading = str(column.get("heading", "")).strip()
            if heading:
                add_textbox(
                    slide,
                    heading,
                    x,
                    1.45,
                    5.5,
                    0.5,
                    size=18,
                    color=theme["accent"],
                    bold=True,
                )
            add_bullets(
                slide,
                column.get("bullets", []),
                x,
                2.0,
                5.45,
                4.45,
                theme,
            )
    else:
        add_bullets(slide, value.get("bullets", []), 0.9, 1.5, 11.5, 4.95, theme)
    add_footer(slide, len(presentation.slides), theme)


def create_section_slide(
    presentation: PresentationType,
    value: dict[str, object],
    theme: dict[str, str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, theme["accent"])
    add_textbox(
        slide,
        str(value.get("title", "")),
        1.0,
        2.4,
        11.2,
        1.0,
        size=34,
        color="FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    subtitle = str(value.get("subtitle", "")).strip()
    if subtitle:
        add_textbox(
            slide,
            subtitle,
            1.3,
            3.55,
            10.6,
            0.65,
            size=18,
            color="E2E8F0",
            align=PP_ALIGN.CENTER,
        )


def command_create(args: argparse.Namespace) -> None:
    spec = read_spec(args.spec)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    theme = theme_from(spec)
    slides = spec.get("slides", [])
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides must be a non-empty array")
    for value in slides:
        if not isinstance(value, dict):
            raise TypeError("every slide must be an object")
        layout = str(value.get("layout", "content")).lower()
        if layout == "title":
            create_title_slide(presentation, value, theme)
        elif layout == "section":
            create_section_slide(presentation, value, theme)
        elif layout in {"content", "two-column"}:
            create_content_slide(presentation, value, theme)
        else:
            raise ValueError(f"unsupported slide layout: {layout}")
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    emit(
        {
            "status": "created",
            "path": str(output.resolve()),
            "slides": len(presentation.slides),
        }
    )


def shape_text(shape: object) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text
    if getattr(shape, "has_table", False):
        return "\n".join(
            "\t".join(cell.text for cell in row.cells) for row in shape.table.rows
        )
    return ""


def command_inspect(args: argparse.Namespace) -> None:
    source = Path(args.input)
    presentation = Presentation(source)
    payload = {
        "path": str(source.resolve()),
        "sizeInches": {
            "width": round(presentation.slide_width / 914400, 3),
            "height": round(presentation.slide_height / 914400, 3),
        },
        "slides": [
            {
                "index": slide_index,
                "shapes": [
                    {
                        "index": shape_index,
                        "name": shape.name,
                        "type": str(shape.shape_type),
                        "text": shape_text(shape),
                        "boundsInches": {
                            "x": round(shape.left / 914400, 3),
                            "y": round(shape.top / 914400, 3),
                            "width": round(shape.width / 914400, 3),
                            "height": round(shape.height / 914400, 3),
                        },
                    }
                    for shape_index, shape in enumerate(slide.shapes)
                ],
            }
            for slide_index, slide in enumerate(presentation.slides, start=1)
        ],
    }
    emit(payload, args.output)


def replace_text_frame(frame: object, old: str, new: str) -> int:
    count = 0
    for paragraph in frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        matches = text.count(old)
        if matches == 0:
            continue
        replacement = text.replace(old, new)
        if paragraph.runs:
            paragraph.runs[0].text = replacement
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = replacement
        count += matches
    return count


def command_replace(args: argparse.Namespace) -> None:
    presentation = Presentation(args.input)
    replacements = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replacements += replace_text_frame(shape.text_frame, args.old, args.new)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replacements += replace_text_frame(
                            cell.text_frame, args.old, args.new
                        )
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    emit(
        {
            "status": "updated",
            "path": str(output.resolve()),
            "replacements": replacements,
        }
    )


def dependency_roots() -> list[Path]:
    configured = os.environ.get("CODEX_WORKSPACE_DEPENDENCIES", "").strip()
    if not configured:
        return []
    root = Path(configured).expanduser()
    return [root, root / "dependencies"]


def find_binary(name: str) -> str | None:
    resolved = shutil.which(name)
    if resolved:
        return resolved
    for root in dependency_roots():
        for candidate in (
            root / "bin" / "override" / name,
            root / "bin" / "fallback" / name,
            root / "native" / "poppler" / "bin" / name,
            root / "native" / "poppler" / "poppler" / "bin" / name,
        ):
            if candidate.is_file():
                return str(candidate)
    return None


def libreoffice_environment() -> dict[str, str]:
    environment = os.environ.copy()
    if environment.get("FONTCONFIG_FILE"):
        return environment
    for root in dependency_roots():
        configs = sorted(
            root.glob("native/libreoffice-headless/**/Resources/fontconfig/fonts.conf")
        )
        if configs:
            environment["FONTCONFIG_FILE"] = str(configs[0])
            environment["FONTCONFIG_PATH"] = str(configs[0].parent)
            break
    return environment


def command_render(args: argparse.Namespace) -> None:
    source = Path(args.input).resolve()
    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    office = find_binary("soffice") or find_binary("libreoffice")
    if not office:
        raise RuntimeError("LibreOffice is required to render PPTX files")
    subprocess.run(
        [
            office,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(source),
        ],
        check=True,
        capture_output=True,
        text=True,
        env=libreoffice_environment(),
    )
    pdf_path = output_dir / f"{source.stem}.pdf"
    if not pdf_path.is_file():
        raise RuntimeError("LibreOffice did not produce the expected PDF")
    images: list[str] = []
    if args.images:
        pdftoppm = find_binary("pdftoppm")
        if not pdftoppm:
            raise RuntimeError("pdftoppm is required when --images is requested")
        prefix = output_dir / source.stem
        subprocess.run(
            [pdftoppm, "-png", "-r", str(args.dpi), str(pdf_path), str(prefix)],
            check=True,
        )
        images = [str(path) for path in sorted(output_dir.glob(f"{source.stem}-*.png"))]
    emit({"pdf": str(pdf_path), "images": images})


def command_validate(args: argparse.Namespace) -> None:
    source = Path(args.input)
    issues: list[str] = []
    warnings: list[str] = []
    try:
        with zipfile.ZipFile(source) as archive:
            corrupt = archive.testzip()
            if corrupt:
                issues.append(f"corrupt ZIP entry: {corrupt}")
            if "ppt/presentation.xml" not in archive.namelist():
                issues.append("missing PPTX presentation part")
        presentation = Presentation(source)
        if not presentation.slides:
            issues.append("presentation has no slides")
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if shape.left < 0 or shape.top < 0:
                    issues.append(
                        f"slide {slide_index}: {shape.name} starts outside the canvas"
                    )
                if shape.left + shape.width > presentation.slide_width:
                    issues.append(
                        f"slide {slide_index}: {shape.name} exceeds the right edge"
                    )
                if shape.top + shape.height > presentation.slide_height:
                    issues.append(
                        f"slide {slide_index}: {shape.name} exceeds the bottom edge"
                    )
                text = shape_text(shape)
                if len(text) > 1200:
                    warnings.append(
                        f"slide {slide_index}: {shape.name} contains dense text"
                    )
    except Exception as error:  # noqa: BLE001 - validation reports format failures
        issues.append(str(error))
    emit(
        {
            "valid": not issues,
            "issues": issues,
            "warnings": warnings,
            "path": str(source.resolve()),
        }
    )
    if issues:
        raise SystemExit(1)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)
    create = subparsers.add_parser("create")
    create.add_argument("--spec", required=True)
    create.add_argument("--output", required=True)
    create.set_defaults(handler=command_create)
    inspect = subparsers.add_parser("inspect")
    inspect.add_argument("--input", required=True)
    inspect.add_argument("--output")
    inspect.set_defaults(handler=command_inspect)
    replace = subparsers.add_parser("replace")
    replace.add_argument("--input", required=True)
    replace.add_argument("--old", required=True)
    replace.add_argument("--new", required=True)
    replace.add_argument("--output", required=True)
    replace.set_defaults(handler=command_replace)
    render = subparsers.add_parser("render")
    render.add_argument("--input", required=True)
    render.add_argument("--output-dir", required=True)
    render.add_argument("--images", action="store_true")
    render.add_argument("--dpi", type=int, default=144)
    render.set_defaults(handler=command_render)
    validate = subparsers.add_parser("validate")
    validate.add_argument("--input", required=True)
    validate.set_defaults(handler=command_validate)
    return parser


def main() -> None:
    args = build_parser().parse_args()
    try:
        args.handler(args)
    except (OSError, ValueError, RuntimeError, subprocess.CalledProcessError) as error:
        print(f"presentations: {error}", file=sys.stderr)
        raise SystemExit(2) from error


if __name__ == "__main__":
    main()
